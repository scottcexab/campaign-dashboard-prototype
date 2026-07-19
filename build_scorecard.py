#!/usr/bin/env python3
"""Build the Insider Threat Campaign executive scorecard PPTX (slide 1)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION

# ---- Palette (Exabeam brand) ----
INK        = RGBColor(0x0A, 0x0A, 0x0A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PAPER      = RGBColor(0xF5, 0xF5, 0xF5)
LINE       = RGBColor(0xE0, 0xE0, 0xE0)
MUTED      = RGBColor(0x70, 0x6E, 0x6B)
BLUE       = RGBColor(0x00, 0x6B, 0xFF)
BLUE_DARK  = RGBColor(0x00, 0x3F, 0xCC)
GREEN      = RGBColor(0x00, 0x9D, 0x00)
GREEN_DK   = RGBColor(0x10, 0x6D, 0x00)
GREEN_BG   = RGBColor(0xEA, 0xF6, 0xEE)
AMBER      = RGBColor(0xB5, 0x7A, 0x00)
AMBER_BG   = RGBColor(0xFD, 0xF2, 0xDD)
BLUE_BG    = RGBColor(0xEA, 0xF2, 0xFC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    set_fill(shp, color)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             font='Arial', anchor=MSO_ANCHOR.TOP, spacing=None, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if spacing:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(spacing))
    return box

def add_pill(slide, x, y, w, h, text, bg, fg):
    shp = add_rect(slide, x, y, w, h, bg, radius=0.5)
    tf = shp.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.name = 'Arial'
    r.font.color.rgb = fg
    return shp

slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)

# ---- Header band ----
add_rect(slide, 0, 0, 13.333, 1.15, INK)
add_text(slide, 0.55, 0.18, 9.5, 0.5, "INSIDER THREAT CAMPAIGN", 22, WHITE, bold=True, font='Arial', spacing=40)
add_text(slide, 0.55, 0.66, 9.5, 0.35, "Executive Scorecard  ·  Jul 1 – Sep 30, 2026", 12.5, RGBColor(0xC9,0xC8,0xDE))
add_text(slide, 9.9, 0.30, 2.9, 0.3, "Prepared: Oct 1, 2026", 10, RGBColor(0x9C,0x9B,0xB4), align=PP_ALIGN.RIGHT)
add_text(slide, 9.9, 0.58, 2.9, 0.3, "Marketing Analytics", 10, RGBColor(0x9C,0x9B,0xB4), align=PP_ALIGN.RIGHT)

# ---- Overall status strip ----
add_rect(slide, 0.55, 1.35, 12.23, 0.62, GREEN_BG, radius=0.25)
add_text(slide, 0.85, 1.35, 4.2, 0.62, "●  ON TRACK", 14, GREEN_DK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_text(slide, 5.0, 1.35, 7.5, 0.62,
         "4 of 6 KPIs met or exceeded goal  ·  Pipeline +37% vs. target  ·  Opportunity creation trailing goal by 14%",
         12, INK, anchor=MSO_ANCHOR.MIDDLE)

# ---- KPI cards ----
kpis = [
    ("Spend vs. Budget",  "$184.3K", "of $250K budget", "73.7%", "ON TRACK", BLUE_BG, BLUE_DARK),
    ("Leads Generated",   "3,842",   "Goal: 3,500",      "110%",  "EXCEEDED", GREEN_BG, GREEN_DK),
    ("MQLs",              "1,120",   "Goal: 1,000",      "112%",  "EXCEEDED", GREEN_BG, GREEN_DK),
    ("Opportunities",     "214",     "Goal: 250",        "86%",   "AT RISK",  AMBER_BG, AMBER),
    ("Pipeline Generated","$4.8M",   "Goal: $3.5M",      "137%",  "EXCEEDED", GREEN_BG, GREEN_DK),
    ("Campaign ROI",      "551%",    "Goal: 400%",       "138%",  "EXCEEDED", GREEN_BG, GREEN_DK),
]
card_w = 1.93
gap = 0.12
start_x = 0.55
card_y = 2.15
card_h = 1.95
for i, (label, val, sub, pct, status, sbg, sfg) in enumerate(kpis):
    x = start_x + i*(card_w+gap)
    add_rect(slide, x, card_y, card_w, card_h, PAPER, radius=0.06)
    add_text(slide, x+0.14, card_y+0.14, card_w-0.28, 0.4, label.upper(), 9, MUTED, bold=True, spacing=10)
    add_text(slide, x+0.14, card_y+0.46, card_w-0.28, 0.55, val, 24, INK, bold=True)
    add_text(slide, x+0.14, card_y+0.98, card_w-0.28, 0.3, sub, 9.5, MUTED)
    add_pill(slide, x+0.14, card_y+1.38, card_w-0.28, 0.32, status + "  ·  " + pct, sbg, sfg)

# ---- Key Takeaways (left) ----
tak_y = 4.20
add_text(slide, 0.55, tak_y, 6.9, 0.35, "KEY TAKEAWAYS", 12, INK, bold=True, spacing=15)
add_rect(slide, 0.55, tak_y+0.38, 0.5, 0.03, BLUE)
takeaways = [
    "Lead and pipeline generation significantly exceeded targets (pipeline +37% vs. goal), driven by Paid Search, Email Nurture, and the new Physical Events channel.",
    "Opportunity creation is trailing goal (214 vs. 250) despite strong top-of-funnel volume — signals an MQL-to-SQL handoff gap worth reviewing with Sales.",
    "PathFactory content scoring directly influenced 148 opportunities (~$1.1M pipeline), validating the TOFU → MOFU → BOFU nurture model.",
    "Campaign is under budget at 73.7% spend utilization with one month remaining — room to reinvest in top-performing tactics (Whitepaper, Webinar Series).",
]
ty = tak_y + 0.55
for t in takeaways:
    add_rect(slide, 0.55, ty+0.09, 0.09, 0.09, BLUE)
    add_text(slide, 0.80, ty, 6.6, 0.5, t, 10.5, RGBColor(0x3E,0x3E,0x3C))
    ty += 0.54

# ---- Right column: pipeline by tactic mini ranking ----
rx = 7.85
add_text(slide, rx, tak_y, 4.9, 0.35, "TOP TACTICS BY PIPELINE", 12, INK, bold=True, spacing=15)
add_rect(slide, rx, tak_y+0.38, 0.5, 0.03, BLUE)
tactics = [
    ("Gated Whitepaper", 1240000, "#0d366b"),
    ("Webinar Series", 980000, "#104281"),
    ("ABM Direct Mail", 760000, "#184f95"),
    ("LinkedIn InMail", 610000, "#1c5cab"),
    ("Exec Roundtables", 520000, "#2a78d6"),
]
max_v = tactics[0][1]
by = tak_y + 0.62
bar_area_w = 3.55
label_w = 1.65
for name, val, hexcolor in tactics:
    rgb = RGBColor(int(hexcolor[1:3],16), int(hexcolor[3:5],16), int(hexcolor[5:7],16))
    add_text(slide, rx, by, label_w, 0.28, name, 9.5, RGBColor(0x3E,0x3E,0x3C))
    bw = max(0.15, (val/max_v) * bar_area_w)
    add_rect(slide, rx+label_w+0.1, by+0.02, bw, 0.24, rgb, radius=0.3)
    add_text(slide, rx+label_w+0.1+bw+0.08, by, 1.3, 0.28, "${:.2f}M".format(val/1_000_000), 9.5, INK, bold=True)
    by += 0.42

# ---- Footer ----
add_rect(slide, 0, 7.18, 13.333, 0.32, PAPER)
add_text(slide, 0.55, 7.20, 9, 0.28,
         "Prototype scorecard — sample data for illustration only. Not a real Exabeam campaign report.",
         8.5, MUTED, italic=True)
add_text(slide, 10.5, 7.20, 2.3, 0.28, "Page 1 of 2", 8.5, MUTED, align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 2 — Campaign Performance Tracking
# Layout inspired by a "sales play performance" template structure:
# header (plan assumptions / plan targets table / trending key) plus
# three body cards (Leads, Opportunity Pipeline, Conversion). Rebuilt
# from scratch with Exabeam branding and Insider Threat Campaign data —
# no branding, text, or assets carried over from the reference template.
# =====================================================================

def style_chart_text(chart, size=9, color=RGBColor(0x3E,0x3E,0x3C)):
    try:
        chart.font.size = Pt(size)
        chart.font.color.rgb = color
        chart.font.name = 'Arial'
    except Exception:
        pass

slide2 = prs.slides.add_slide(blank)
add_rect(slide2, 0, 0, 13.333, 7.5, WHITE)

# ---- Header band ----
add_rect(slide2, 0, 0, 13.333, 0.85, INK)
add_text(slide2, 0.55, 0.14, 9.5, 0.4, "INSIDER THREAT CAMPAIGN — PERFORMANCE TRACKING", 17, WHITE, bold=True, spacing=30)
add_text(slide2, 0.55, 0.50, 9.5, 0.3, "Source: Campaign Dashboard  ·  Data as of Oct 1, 2026", 10.5, RGBColor(0xC9,0xC8,0xDE))

# ---- Row 1: Plan Assumptions / Plan Targets table / Trending Key ----
row1_y = 1.02
row1_h = 1.35

# Plan Assumptions
add_rect(slide2, 0.4, row1_y, 2.75, row1_h, PAPER, radius=0.06)
add_text(slide2, 0.58, row1_y+0.12, 2.4, 0.28, "PLAN ASSUMPTIONS", 9.5, MUTED, bold=True, spacing=10)
assumptions = ["Avg Deal Size: ~$22.4K", "Lead → MQL: 29.2%", "MQL → Opportunity: 19.1%"]
ay = row1_y + 0.44
for a in assumptions:
    add_text(slide2, 0.58, ay, 2.4, 0.26, a, 10.5, INK)
    ay += 0.28

# Plan Targets table
tbl_x, tbl_w = 3.35, 5.9
add_text(slide2, tbl_x, row1_y-0.02, 3, 0.24, "PLAN TARGETS", 9.5, MUTED, bold=True, spacing=10)
gfx = slide2.shapes.add_table(5, 4, Inches(tbl_x), Inches(row1_y+0.20), Inches(tbl_w), Inches(row1_h-0.20))
tbl = gfx.table
tbl.columns[0].width = Inches(1.7)
tbl.columns[1].width = Inches(1.4)
tbl.columns[2].width = Inches(1.4)
tbl.columns[3].width = Inches(1.4)
headers = ["Metric", "Goal", "Actual", "% of Goal"]
rows_data = [
    ("Leads",         "3,500",  "3,842", "110%"),
    ("MQLs",          "1,000",  "1,120", "112%"),
    ("Opportunities", "250",    "214",   "86%"),
    ("Pipeline",      "$3.5M",  "$4.8M", "137%"),
]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = INK
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci==0 else PP_ALIGN.RIGHT
    r = p.runs[0]; r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name='Arial'
    cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(6); cell.margin_right = Pt(6)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
for ri, row in enumerate(rows_data, start=1):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = PAPER if ri % 2 == 0 else WHITE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci==0 else PP_ALIGN.RIGHT
        r = p.runs[0]; r.font.size = Pt(9.5); r.font.bold = (ci==0); r.font.color.rgb = INK; r.font.name='Arial'
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(6); cell.margin_right = Pt(6)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# Trending Key legend
key_x = 9.55
add_rect(slide2, key_x, row1_y, 3.38, row1_h, PAPER, radius=0.06)
add_text(slide2, key_x+0.18, row1_y+0.12, 3, 0.24, "TRENDING KEY", 9.5, MUTED, bold=True, spacing=10)
trend_items = [("On Plan", GREEN_DK, GREEN_BG), ("Needs Attention", AMBER, AMBER_BG), ("Under Performing", RGBColor(0xC2,0x39,0x34), RGBColor(0xFC,0xF2,0xF2))]
ky = row1_y + 0.46
for label, fg, bg in trend_items:
    dot = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(key_x+0.18), Inches(ky+0.03), Inches(0.14), Inches(0.14))
    set_fill(dot, fg); dot.shadow.inherit = False
    add_text(slide2, key_x+0.42, ky, 2.8, 0.26, label, 10.5, INK)
    ky += 0.29

# ---- Row 2: three body cards ----
row2_y = 2.55
row2_h = 4.35
card_gap = 0.18
card3_w = (13.333 - 0.8 - 2*card_gap) / 3
xs = [0.4, 0.4 + card3_w + card_gap, 0.4 + 2*(card3_w + card_gap)]

# ---- Card 1: Leads ----
cx = xs[0]
add_rect(slide2, cx, row2_y, card3_w, row2_h, PAPER, radius=0.04)
add_rect(slide2, cx, row2_y, card3_w, 0.42, BLUE)
add_text(slide2, cx+0.16, row2_y, card3_w-0.32, 0.42, "LEADS", 11, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, spacing=15)
add_text(slide2, cx+0.16, row2_y+0.55, card3_w-0.32, 0.3, "Total Leads: 3,842", 11.5, INK, bold=True)
add_text(slide2, cx+0.16, row2_y+0.83, card3_w-0.32, 0.26, "Avg Weekly Growth: +9.1%", 10, MUTED)

leads_chart_data = CategoryChartData()
leads_chart_data.categories = ['W1','W2','W3','W4','W5','W6','W7','W8','W9','W10','W11','W12','W13']
leads_chart_data.add_series('Leads', (140,205,260,310,295,340,380,410,455,420,470,510,490))
gframe = slide2.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(cx+0.1), Inches(row2_y+1.2), Inches(card3_w-0.2), Inches(2.0), leads_chart_data)
lc = gframe.chart
lc.has_legend = False
lc.has_title = False
plot = lc.plots[0]
series = plot.series[0]
series.format.line.color.rgb = BLUE
series.format.line.width = Pt(2.25)
series.smooth = False
try:
    series.marker.format.fill.solid(); series.marker.format.fill.fore_color.rgb = BLUE
    series.marker.size = 5
except Exception:
    pass
cat_ax = lc.category_axis; val_ax = lc.value_axis
cat_ax.tick_labels.font.size = Pt(7.5); val_ax.tick_labels.font.size = Pt(7.5)
cat_ax.format.line.color.rgb = LINE; val_ax.format.line.color.rgb = LINE
val_ax.has_major_gridlines = True
val_ax.major_gridlines.format.line.color.rgb = LINE
val_ax.major_gridlines.format.line.width = Pt(0.5)

add_text(slide2, cx+0.16, row2_y+3.32, card3_w-0.32, 0.22, "TOP 3 CHANNELS", 9, MUTED, bold=True, spacing=10)
top_channels = [("Paid Search", "1,040"), ("Email Nurture", "980"), ("Webinar", "612")]
tcy = row2_y + 3.56
for name, leads_n in top_channels:
    add_text(slide2, cx+0.16, tcy, card3_w-1.1, 0.22, name, 9.5, INK)
    add_text(slide2, cx+card3_w-1.0, tcy, 0.84, 0.22, leads_n, 9.5, INK, bold=True, align=PP_ALIGN.RIGHT)
    tcy += 0.225

# ---- Card 2: Opportunity Pipeline ----
cx = xs[1]
add_rect(slide2, cx, row2_y, card3_w, row2_h, PAPER, radius=0.04)
add_rect(slide2, cx, row2_y, card3_w, 0.42, GREEN)
add_text(slide2, cx+0.16, row2_y, card3_w-0.32, 0.42, "OPPORTUNITY PIPELINE", 11, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, spacing=15)
add_text(slide2, cx+0.16, row2_y+0.55, card3_w-0.32, 0.3, "YTD/QTD Pipeline: $4.8M", 11.5, INK, bold=True)
add_text(slide2, cx+0.16, row2_y+0.83, card3_w-0.32, 0.26, "Avg Sales Cycle: 46 days", 10, MUTED)

funnel_chart_data = CategoryChartData()
funnel_chart_data.categories = ['Inquiry','MQL','SQL','Opportunity','Won']
funnel_chart_data.add_series('Count', (3842,1120,412,214,38))
gframe2 = slide2.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(cx+0.1), Inches(row2_y+1.2), Inches(card3_w-0.2), Inches(2.15), funnel_chart_data)
bc = gframe2.chart
bc.has_legend = False
bc.has_title = False
bplot = bc.plots[0]
bplot.gap_width = 40
bseries = bplot.series[0]
bseries.format.fill.solid(); bseries.format.fill.fore_color.rgb = GREEN
bcat_ax = bc.category_axis; bval_ax = bc.value_axis
bcat_ax.tick_labels.font.size = Pt(8); bval_ax.tick_labels.font.size = Pt(7.5)
bcat_ax.format.line.color.rgb = LINE; bval_ax.format.line.color.rgb = LINE
bval_ax.has_major_gridlines = True
bval_ax.major_gridlines.format.line.color.rgb = LINE
bval_ax.major_gridlines.format.line.width = Pt(0.5)

add_text(slide2, cx+0.16, row2_y+3.5, card3_w-0.32, 0.24, "CLOSED WON / LOST", 9, MUTED, bold=True, spacing=10)
add_text(slide2, cx+0.16, row2_y+3.75, card3_w-0.32, 0.28, "38 Won  /  14 Lost  (73% win rate)", 10.5, INK)

# ---- Card 3: Conversion & Key Opportunities ----
cx = xs[2]
add_rect(slide2, cx, row2_y, card3_w, row2_h, PAPER, radius=0.04)
add_rect(slide2, cx, row2_y, card3_w, 0.42, RGBColor(0x8D,0x00,0xFF))
add_text(slide2, cx+0.16, row2_y, card3_w-0.32, 0.42, "CONVERSION TRENDS", 11, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, spacing=15)

conv_y = row2_y + 0.55
add_text(slide2, cx+0.16, conv_y, 1.7, 0.26, "MQLs MoM", 9.5, MUTED)
add_text(slide2, cx+0.16, conv_y+0.24, 1.7, 0.32, "+9.7%", 15, GREEN_DK, bold=True)
add_rect(slide2, cx+card3_w/2, conv_y+0.02, 0.012, 0.5, LINE)
add_text(slide2, cx+card3_w/2+0.16, conv_y, 1.7, 0.26, "Opps MoM", 9.5, MUTED)
add_text(slide2, cx+card3_w/2+0.16, conv_y+0.24, 1.7, 0.32, "-3.1%", 15, AMBER, bold=True)

add_text(slide2, cx+0.16, row2_y+1.25, card3_w-0.32, 0.24, "KEY OPPORTUNITIES", 9.5, MUTED, bold=True, spacing=10)
key_opps = [("Global Bank Co.", "$340K"), ("Meridian Health", "$280K"), ("Northline Retail", "$210K")]
koy = row2_y + 1.52
for name, amt in key_opps:
    add_rect(slide2, cx+0.16, koy, card3_w-0.32, 0.44, WHITE, radius=0.15)
    add_text(slide2, cx+0.30, koy+0.09, card3_w-0.9, 0.26, name, 10.5, INK, bold=True)
    add_text(slide2, cx+card3_w-1.0, koy+0.09, 0.75, 0.26, amt, 10.5, GREEN_DK, bold=True, align=PP_ALIGN.RIGHT)
    koy += 0.52

add_text(slide2, cx+0.16, koy+0.06, card3_w-0.32, 0.24, "TOP REGIONS BY PIPELINE", 9.5, MUTED, bold=True, spacing=10)
regions_mini = [("North America", 2160000, BLUE), ("EMEA", 912000, GREEN), ("APAC", 576000, AMBER)]
max_r = regions_mini[0][1]
ry = koy + 0.32
bar_area_w2 = card3_w - 1.9
label_w2 = 1.05
for name, val, rgb in regions_mini:
    add_text(slide2, cx+0.16, ry, label_w2, 0.24, name, 9, RGBColor(0x3E,0x3E,0x3C))
    bw2 = max(0.1, (val/max_r) * bar_area_w2)
    add_rect(slide2, cx+0.16+label_w2, ry+0.02, bw2, 0.18, rgb, radius=0.3)
    add_text(slide2, cx+0.16+label_w2+bw2+0.06, ry-0.01, 0.9, 0.24, "${:.2f}M".format(val/1_000_000), 8.5, INK, bold=True)
    ry += 0.30

# ---- Footer ----
add_rect(slide2, 0, 7.18, 13.333, 0.32, PAPER)
add_text(slide2, 0.55, 7.20, 9, 0.28,
         "Prototype tracking slide — sample data for illustration only. Layout structure adapted from an internal template; no external branding or content included.",
         8, MUTED, italic=True)
add_text(slide2, 10.5, 7.20, 2.3, 0.28, "Page 2 of 2", 8.5, MUTED, align=PP_ALIGN.RIGHT)

prs.save('/private/tmp/claude-501/-Users-scottclinton/b975b1c5-0616-4716-ad2d-ab6bcf229c90/scratchpad/campaign-dashboard/Insider_Threat_Campaign_Scorecard.pptx')
print("Saved.")
